import os
import re
import json
import logging
import base64
import openai
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image


# Set up logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
    filename="ai_presentation_generator.log"
)

# Constants for PPT styling
FONT_NAME = "Mangal"
HEADER_FONT_SIZE = Pt(28)
CONTENT_FONT_SIZE = Pt(28)
HEADER_TEXT_WHITE = RGBColor(255, 255, 255)
BACKGROUND_COLOR = RGBColor(0, 0, 0)
HEADER_RED = RGBColor(252, 5, 5)
HEADER_YELLOW = RGBColor(255, 255, 0)
IMAGE1_PATH = "logo.jpg"  # Change this to your actual logo image path
HEADER_RCB_RED = RGBColor(252, 5, 5)
MAX_TEXT_WIDTH = 700  # Define maximum width for text wrapping
LINE_SPACING = Pt(8)

class AIPresentationGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(14)
        self.prs.slide_height = Inches(7.5)

    def add_navigation_bar(self, slide, title, topic, teacher_name):
        """Add the navigation bar with title, logo, and name."""
        # Title Section
        nav_left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(0), Inches(8), Inches(0.5))
        nav_left.fill.solid()
        nav_left.fill.fore_color.rgb = HEADER_RCB_RED
        nav_left.text_frame.text = title
        nav_left.text_frame.paragraphs[0].font.size = Pt(20)
        nav_left.text_frame.paragraphs[0].font.bold = True
        nav_left.text_frame.paragraphs[0].font.color.rgb = HEADER_TEXT_WHITE
        nav_left.text_frame.paragraphs[0].font.name = FONT_NAME
        nav_left.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Add logo
        slide.shapes.add_picture(IMAGE1_PATH, Inches(0), Inches(0), Inches(1), Inches(0.5))

        # Middle Section for Topic
        nav_middle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(0), Inches(6), Inches(0.5))
        nav_middle.fill.solid()
        nav_middle.fill.fore_color.rgb = HEADER_YELLOW
        nav_middle.text_frame.clear()
        nav_middle.text_frame.text = topic
        nav_middle.text_frame.paragraphs[0].font.size = Pt(20)
        nav_middle.text_frame.paragraphs[0].font.bold = True
        nav_middle.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        nav_middle.text_frame.paragraphs[0].font.name = FONT_NAME
        nav_middle.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Right Section for Creator's Name
        nav_right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10), Inches(0), Inches(4), Inches(0.5))
        nav_right.fill.solid()
        nav_right.fill.fore_color.rgb = HEADER_RCB_RED
        nav_right.text_frame.text = f"BY: {teacher_name}"
        nav_right.text_frame.paragraphs[0].font.size = Pt(20)
        nav_right.text_frame.paragraphs[0].font.bold = True
        nav_right.text_frame.paragraphs[0].font.color.rgb = HEADER_TEXT_WHITE
        nav_right.text_frame.paragraphs[0].font.name = FONT_NAME
        nav_right.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    def wrap_text_to_fit(self, text, max_width=MAX_TEXT_WIDTH):
        """Wrap text to fit within specified width."""
        words = text.split()
        lines = []
        current_line = []
        current_length = 0
        
        for word in words:
            word_length = len(word)
            if current_length + word_length <= 60:  # Approximate characters per line
                current_line.append(word)
                current_length += word_length + 1
            else:
                lines.append(' '.join(current_line))
                current_line = [word]
                current_length = word_length + 1
        
        if current_line:
            lines.append(' '.join(current_line))
        
        return '\n'.join(lines)

    def add_main_title(self, slide, title_text):
        """Adds a title to the slide."""
        title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.75), Inches(12), Inches(0.8))
        text_frame = title_shape.text_frame
        text_frame.text = title_text
        text_frame.paragraphs[0].font.size = Pt(32)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        text_frame.paragraphs[0].font.name = FONT_NAME
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def add_content_textbox(self, slide):
        """Adds a textbox for bullet points and returns its text frame."""
        top_margin = Inches(1.75)  # Adjusted for title and nav bar
        bottom_margin = Inches(0.5)
        side_margin = Inches(1)
        
        max_textbox_height = Inches(5.5)  # Allows more content before requiring a new slide
    
        textbox = slide.shapes.add_textbox(
            left=side_margin,
            top=top_margin,
            width=Inches(12),  # 14 inches (slide width) - 2 inches (margins)
            height=max_textbox_height  # Adjusted for better fitting
        )
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        return text_frame

    def add_point_to_textbox(self, text_frame, point,available_height):
        """Adds a bullet point to the given text frame while checking if it fits."""
        wrapped_text, overflow_text = self.split_text_to_fit(text_frame, point, available_height)
        
        p = text_frame.add_paragraph()
        p.text = f"➤ {wrapped_text}"  # ✅ Set text first before applying font styles
        p.font.size = CONTENT_FONT_SIZE
        p.space_after = Pt(12)  # Consistent spacing between points
        p.font.name = FONT_NAME
        p.font.color.rgb = RGBColor(255, 255, 0)

        return overflow_text  # ✅ Return any text that didn't fit

    # def create_presentation(self, extracted_data, title, topic, teacher_name):
    #     """Generate a PowerPoint presentation with optimal text splitting and full slide utilization."""
    #     print("📢 Generating Fully Optimized PPT...")

    #     prs = Presentation()
    #     prs.slide_width = Inches(14)
    #     prs.slide_height = Inches(7.5)

    #     slide_count = 0
    #     max_textbox_height = Inches(5)  # Max content height per slide

    #     current_slide = None
    #     current_text_frame = None
    #     available_height = max_textbox_height  # Reset available space per slide

    #     print(f"📢 Received Extracted Data Type: {type(extracted_data)} - {extracted_data}")

    #     # ✅ Check if extracted_data is a string
    #     if isinstance(extracted_data, str):
    #         extracted_data = json.loads(extracted_data)  # Convert to list

    #     if not isinstance(extracted_data, list):
    #         print(f"❌ Invalid extracted_data format! Expected a list but got: {type(extracted_data)}")
    #         return None 

    #     for i, content in enumerate(extracted_data):
    #         title_text = content.get("title", "")
    #         points = content.get("points", [])

    #         if not points:
    #             print(f"⚠️ No points found for slide {i + 1}, skipping.")
    #             continue

    #         for point in points:
    #             wrapped_text = self.wrap_text_to_fit(point)

    #             # Check if there's enough space to fit the next point, else create a new slide
    #             if available_height < self.estimate_text_height(wrapped_text):
    #                 current_slide = prs.slides.add_slide(prs.slide_layouts[6])
    #                 current_slide.background.fill.solid()
    #                 current_slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)

    #                 # Add navigation and title
    #                 self.add_navigation_bar(current_slide, title, topic, teacher_name)
    #                 if title_text:
    #                     self.add_main_title(current_slide, title_text)

    #                 # Add content box
    #                 current_text_frame = self.add_content_textbox(current_slide)
    #                 available_height = max_textbox_height  # Reset available height
    #                 slide_count+=1        

    #             # Now add the point to the slide
    #             self.add_point_to_textbox(current_text_frame, wrapped_text)
    #             available_height -= self.estimate_text_height(wrapped_text) + Inches(0.2)  # Reduce available space

                  

    #     if slide_count == 0:
    #         print("❌ No valid slides created. Cannot save an empty presentation.")
    #         return None

    #     timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    #     filename = f"final_presentation_{timestamp}.pptx"
    #     prs.save(filename)

    #     print(f"✅ Presentation saved as {filename}")
    #     return filename

    def create_presentation(self, extracted_data, title, topic, teacher_name):
        """Generate a PowerPoint presentation with proper slide utilization."""
        print("📢 Generating Fully Optimized PPT...")

        prs = Presentation()
        prs.slide_width = Inches(14)
        prs.slide_height = Inches(7.5)

        slide_count = 0
        max_textbox_height = Inches(5)  # Max content height per slide

        current_slide = None
        current_text_frame = None
        available_height = max_textbox_height  # Reset available space per slide

        print(f"📢 Received Extracted Data Type: {type(extracted_data)} - {extracted_data}")

        # ✅ Check if extracted_data is a string
        if isinstance(extracted_data, str):
            extracted_data = json.loads(extracted_data)  # Convert to list

        if not isinstance(extracted_data, list):
            print(f"❌ Invalid extracted_data format! Expected a list but got: {type(extracted_data)}")
            return None 

        for i, content in enumerate(extracted_data):
            title_text = content.get("title", "")
            points = content.get("points", [])

            if not points:
                print(f"⚠️ No points found for slide {i + 1}, skipping.")
                continue

            for point in points:
                wrapped_text = self.wrap_text_to_fit(point)
                point_height = self.estimate_text_height(wrapped_text) + Inches(0.2)  # Extra spacing

                # ✅ If no slide exists, create one
                if current_slide is None:
                    current_slide = prs.slides.add_slide(prs.slide_layouts[6])
                    current_slide.background.fill.solid()
                    current_slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)

                    # Add navigation and title
                    self.add_navigation_bar(current_slide, title, topic, teacher_name)
                    if title_text:
                        self.add_main_title(current_slide, title_text)

                    # Add content box
                    current_text_frame = self.add_content_textbox(current_slide)
                    available_height = max_textbox_height  # Reset available height
                    slide_count += 1

                # ✅ If the point fits, add it to the current slide
                if available_height >= point_height:
                    self.add_point_to_textbox(current_text_frame, wrapped_text,available_height)
                    available_height -= point_height  # Reduce available space
                else:
                    # ✅ If it doesn't fit, create a new slide and then add it
                    current_slide = prs.slides.add_slide(prs.slide_layouts[6])
                    current_slide.background.fill.solid()
                    current_slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)

                    # Add navigation and title
                    self.add_navigation_bar(current_slide, title, topic, teacher_name)
                    if title_text:
                        self.add_main_title(current_slide, title_text)

                    # Add content box
                    current_text_frame = self.add_content_textbox(current_slide)
                    available_height = max_textbox_height  # Reset available height
                    slide_count += 1

                    # Now add the point
                    self.add_point_to_textbox(current_text_frame, wrapped_text,available_height)
                    available_height -= point_height  # Reduce available space

        if slide_count == 0:
            print("❌ No valid slides created. Cannot save an empty presentation.")
            return None

        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"final_presentation_{timestamp}.pptx"
        prs.save(filename)

        print(f"✅ Presentation saved as {filename}")
        return filename

    def split_text_to_fit(self, text_frame, text, available_height):
        """Splits text so that as much as possible remains in the slide, and overflow moves to the next slide."""
        words = text.split()
        fitted_text = ""
        remaining_text = ""
        estimated_height = 0
        max_line_length = 50  # Approximate max characters per line

        wrapped_lines = []
        current_line = []

        for word in words:
            if len(" ".join(current_line) + " " + word) <= max_line_length:
                current_line.append(word)
            else:
                wrapped_lines.append(" ".join(current_line))
                current_line = [word]

        if current_line:
            wrapped_lines.append(" ".join(current_line))

        for line in wrapped_lines:
            line_height = Pt(28) / 72  # Approximate line height in inches
            if estimated_height + line_height > available_height:
                remaining_text += line + " "
            else:
                fitted_text += line + " "
                estimated_height += line_height

        return fitted_text.strip(), remaining_text.strip()

    def estimate_text_height(self, text):
        """More accurate estimation of text height in PowerPoint slides."""
        words_per_line = 50  # Approximate max characters per line
        line_height = Pt(28) / 72  # Line height in inches

        # Ensure accurate line count based on real wrapping
        lines = (len(text) // words_per_line) + 1
        return Inches(lines * line_height)


if __name__ == "__main__":
    extracted_data = [
        {
            "title": "अनुपात (Ratio)",
            "points": [
                "अनुपात वह गणितीय छोजक है जो समान इकाई की दो असमान राशियों के बीच तुलना दिखाता है।",
                "Ratio में कभी भी Number से गुणा करने पर कोई फर्क नहीं पड़ता है।",
                "लेकिन जोड़ने और घटाने पर फर्क पड़ता है। जब दो अनुपात आपस में बराबर हो तो इसे समानुपात कहते ",
                "लेकिन जोड़ने और घटाने पर फर्क पड़ता है। जब दो अनुपात आपस में बराबर हो तो इसे समानुपात कहते हैं।",
                "इसे '=' चिह्न से दर्शाया जाता है।",
                "दो या दो से अधिक अनुपात के अंतिम पदों के गुणनफल को जटिल अनुपात कहते हैं।",
                "जैसे कि a:b, c:d, e:f को गुणा करने पर (a×c×e):(b×d×f) एक जटिल अनुपात बनता है।"
            ]
        },
        {
            "title": "वर्गमूलानुपात (Sub-Duplicate Ratio)",
            "points": [
                "यदि m:n का वर्गमूल लिया जाए तो उसे वर्गमूलानुपात कहते हैं।"
            ]
        },
        {
            "title": "घनानुपात (Triplicate Ratio)",
            "points": [
                "यदि m:n कोई अनुपात हो, तो m³:n³ को घनानुपात कहते हैं।"
            ]
        }
    ]

    handler = AIPresentationGenerator()
    handler.create_presentation(
        extracted_data=extracted_data,
        title="NEXT LEVEL ACADEMY",
        topic="Ratio",
        teacher_name="R K Gupta"
    )