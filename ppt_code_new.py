from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Constants for styling

HEADER_TEXT_COLOR = RGBColor(255, 255, 255)  # White
ANSWER_TEXT_COLOR = RGBColor(0, 255, 0)  # Green for answers
FONT_NAME = "Calibri"
HEADER_FONT_SIZE = Pt(16)
QUESTION_FONT_SIZE = Pt(20)
OPTION_FONT_SIZE = Pt(18)
ANSWER_FONT_SIZE = Pt(22)
BACKGROUND_COLOR = RGBColor(0, 0, 0)  # Black
HEADER_RCB_RED = RGBColor(252, 5, 5)
HEADER_YELLOW = RGBColor(255, 255, 0)  # Yellow
HEADER_TEXT_WHITE = RGBColor(255, 255, 255)
QUESTION_TEXT_COLOR = RGBColor(255, 255, 0)  # Yellow
OPTION_TEXT_COLOR = RGBColor(255, 255, 255)  # White
FONT_NAME = "Calibri"
HEADER_FONT_SIZE = Pt(16)
MAX_FONT_SIZE = Pt(22)
MIN_FONT_SIZE = Pt(12)


IMAGE_1_PATH = "logo.jpg"
# IMAGE_2_PATH = "testcase1.jpeg"

@staticmethod
def _calculate_dynamic_font_size(text, max_font_size=Pt(36), min_font_size=Pt(16)):
    """
    Dynamically calculate font size based on text length.
    
    :param text: Input text
    :param max_font_size: Maximum font size
    :param min_font_size: Minimum font size
    :return: Appropriate font size
    """
    # Basic length-based font size calculation
    length = len(text)
    
    # Adjust font size based on text length
    if length < 50:
        return max_font_size  # Large font for short questions
    elif length < 100:
        return Pt(30)
    elif length < 200:
        return Pt(24)
    elif length < 300:
        return Pt(20)
    else:
        return min_font_size

@staticmethod
def wrap_text_to_fit(text, font_size=Pt(36), max_width=700):
    """
    Wrap text to fit within specified width, with dynamic font sizing.
    
    :param text: Input text
    :param font_size: Base font size 
    :param max_width: Maximum width in points
    :return: Wrapped text with adjusted dimensions
    """
    # Dynamically adjust font size based on text length
    adjusted_font_size = _calculate_dynamic_font_size(text)
    
    words = text.split()
    lines = []
    current_line = ""
    approx_char_width = adjusted_font_size.pt * 0.4  # Approximation of character width in points
    
    for word in words:
        if (len(current_line) + len(word)) * approx_char_width > max_width:
            lines.append(current_line.strip())
            current_line = word
        else:
            current_line += " " + word if current_line else word
    
    if current_line:
        lines.append(current_line.strip())

    # Calculate line height with some extra spacing
    line_height = adjusted_font_size.pt * 1.2
    total_height = len(lines) * line_height

    # Return wrapped text and adjust height calculation
    return "\n".join(lines), total_height

def add_question_slide(prs, question_number, question, year, options, title, topic, teacher_name):
    """Add a slide for the question with improved sizing and positioning."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Set slide background color
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BACKGROUND_COLOR

    # Add the header with title, logo, and name
    add_navigation_bar(slide, title, topic, teacher_name)

    # Determine question box dimensions dynamically
    font_size = Pt(36)
    wrapped_question, text_height = wrap_text_to_fit(question, font_size=font_size, max_width=700)
    
    # Adjust box dimensions based on text height
    box_width = Inches(6.3)  # Default width for the box
    lines = wrapped_question.count('\n') + 1 
    
    # Calculate dynamic box height with padding
    box_height = Inches((text_height / 72) + 0.2)  # Convert points to inches, with padding

    # Create question box
    question_box = slide.shapes.add_textbox(Inches(7.5), Inches(1), box_width, box_height)

    # Add question number box
    question_number_box = slide.shapes.add_textbox(Inches(6.3), Inches(1.3), Inches(0.7), Inches(0.9))
    question_number_frame = question_number_box.text_frame
    question_number_frame.clear()
    question_number_paragraph = question_number_frame.add_paragraph()
    question_number_paragraph.text = str(question_number) + "."  # Add the question number
    question_number_box.fill.solid()
    question_number_box.fill.fore_color.rgb = RGBColor(255, 0, 0)
    question_number_paragraph.font.bold = True
    question_number_paragraph.font.color.rgb = RGBColor(255, 255, 255)
    question_number_paragraph.font.name = FONT_NAME
    question_number_paragraph.font.size = Pt(30)
    question_number_paragraph.alignment = PP_ALIGN.RIGHT

    # Populate the question text box
    question_frame = question_box.text_frame
    question_frame.clear()
    question_frame.word_wrap = True
    question_frame.margin_top = Pt(0)
    question_frame.margin_bottom = Pt(0)
    question_frame.margin_left = Pt(0)
    question_frame.margin_right = Pt(0)
    question_paragraph = question_frame.add_paragraph()
    question_paragraph.text = wrapped_question
    question_paragraph.font.bold = True
    question_paragraph.font.color.rgb = QUESTION_TEXT_COLOR
    question_paragraph.font.name = FONT_NAME
    
    # Dynamically adjust font size
    dynamic_font_size = _calculate_dynamic_font_size(question)
    question_paragraph.font.size = dynamic_font_size
    question_paragraph.alignment = PP_ALIGN.LEFT

    # Add the year, if available
    total_options_height = 0
    wrapped_options = []
        
        # Process each option for wrapping and height calculation
    for option in options:
        wrapped_option, option_height = wrap_text_to_fit(
            option, 
            font_size=OPTION_FONT_SIZE,
            max_width=600  # Slightly smaller than box width to ensure padding
        )
        wrapped_options.append(wrapped_option)
        total_options_height += option_height + Pt(10).pt  # Add padding between options
    
    # Calculate optimal position for options box
    # Start options 1 inch below question box, or at slide middle, whichever is higher
    min_options_top = question_box.top + question_box.height + Inches(1)
    slide_middle = (prs.slide_height - Inches(1)) / 2  # Leave space for header
    options_top = max(min_options_top, slide_middle - Inches(total_options_height/72/2))
    
    # Create options box with calculated height
    options_box = slide.shapes.add_textbox(
        question_box.left, 
        options_top, 
        question_box.width, 
        Inches(total_options_height/72 + 0.5)  # Convert points to inches, add padding
    )
    
    options_frame = options_box.text_frame
    options_frame.clear()
    options_frame.word_wrap = True
    
    # Add each wrapped option
    for wrapped_option in wrapped_options:
        option_paragraph = options_frame.add_paragraph()
        option_paragraph.text = wrapped_option
        option_paragraph.font.size = OPTION_FONT_SIZE
        option_paragraph.line_spacing = 1.2  # Proportional line spacing
        option_paragraph.space_after = Pt(10)  # Space between options
        option_paragraph.font.color.rgb = OPTION_TEXT_COLOR
        option_paragraph.font.name = FONT_NAME
        option_paragraph.alignment = PP_ALIGN.LEFT # Left-align the options
    


def create_custom_presentation(data,title,topic,teacher_name):
    """Create a custom PowerPoint presentation based on the provided data."""
    prs = Presentation()
    prs.slide_width = Inches(14)
    prs.slide_height = Inches(7.5)

    question_count = len(data["Question"])
    year_count = len(data.get("Year",[]))
    options_count = len(data.get("Options", []))

    for i in range(question_count):
        question = data["Question"][i]
        year=data["Year"][i] if i<year_count else ""
        options = data["Options"][i] if i < options_count else ""


        # Add question slide
        add_question_slide(prs, i + 1, question,year, options,title,topic,teacher_name)
        # Add answer slide
        #add_answer_slide(prs, i + 1, question, answer, explanation)

    # Save the presentation
    timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
    filename = f"ppt_{timestamp}.pptx"
    prs.save(filename)
    print(f"Presentation saved as {filename}")


def add_navigation_bar(slide,title,topic,teacher_name):
    """Add the navigation bar with title, logo, and name."""
    # Title Section
    nav_left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(0), Inches(8), Inches(0.5))
    nav_left.fill.solid()
    nav_left.fill.fore_color.rgb = HEADER_RCB_RED
    nav_left.text_frame.text = title
    nav_left.text_frame.paragraphs[0].font.size = HEADER_FONT_SIZE
    nav_left.text_frame.paragraphs[0].font.bold = True
    nav_left.text_frame.paragraphs[0].font.color.rgb = HEADER_TEXT_WHITE
    nav_left.text_frame.paragraphs[0].font.name = FONT_NAME
    nav_left.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Add logo
    slide.shapes.add_picture(IMAGE_1_PATH, Inches(0), Inches(0), Inches(1), Inches(0.5))

    # Middle Section for Topic
    nav_middle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(0), Inches(6), Inches(0.5))
    nav_middle.fill.solid()
    nav_middle.fill.fore_color.rgb = HEADER_YELLOW
    nav_middle.text_frame.clear()  # Clear any default text
    nav_middle.text_frame.text = topic  # Add your topic here
    nav_middle.text_frame.paragraphs[0].font.size = HEADER_FONT_SIZE
    nav_middle.text_frame.paragraphs[0].font.bold = True
    nav_middle.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    nav_middle.text_frame.paragraphs[0].font.name = FONT_NAME
    nav_middle.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Right Section for Creator's Name
    nav_right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10), Inches(0), Inches(4), Inches(0.5))
    nav_right.fill.solid()
    nav_right.fill.fore_color.rgb = HEADER_RCB_RED
    nav_right.text_frame.text = f"BY: {teacher_name}"
    nav_right.text_frame.paragraphs[0].font.size = HEADER_FONT_SIZE
    nav_right.text_frame.paragraphs[0].font.bold = True
    nav_right.text_frame.paragraphs[0].font.color.rgb = HEADER_TEXT_WHITE
    nav_right.text_frame.paragraphs[0].font.name = FONT_NAME
    nav_right.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_custom_slide(prs, question_number, question, options, year):
    """Add a single custom slide with dynamically allocated spaces."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Set slide background color
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BACKGROUND_COLOR

    # Navigation Bar Sections
    nav_left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(0), Inches(5), Inches(0.5))
    nav_left.fill.solid()
    nav_left.fill.fore_color.rgb = HEADER_RCB_RED
    nav_left.text_frame.text = "NEXT LEVEL ACADEMY"
    nav_left.text_frame.paragraphs[0].font.size = HEADER_FONT_SIZE
    nav_left.text_frame.paragraphs[0].font.bold = True
    nav_left.text_frame.paragraphs[0].font.color.rgb = HEADER_YELLOW
    nav_left.text_frame.paragraphs[0].font.name = FONT_NAME
    nav_left.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    slide.shapes.add_picture(IMAGE_1_PATH, Inches(0), Inches(0), Inches(1), Inches(0.5))

    nav_middle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(0), Inches(5), Inches(0.5))
    nav_middle.fill.solid()
    nav_middle.fill.fore_color.rgb = HEADER_YELLOW
    # nav_middle.text_frame.text = "Indian Geography Introduction of India"
    nav_middle.text_frame.paragraphs[0].font.size = HEADER_FONT_SIZE
    nav_middle.text_frame.paragraphs[0].font.bold = True
    nav_middle.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    nav_middle.text_frame.paragraphs[0].font.name = FONT_NAME
    nav_middle.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # slide.shapes.add_picture(IMAGE_2_PATH, Inches(5.5), Inches(0), Inches(1), Inches(0.5))

    nav_right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11), Inches(0), Inches(3), Inches(0.5))
    nav_right.fill.solid()
    nav_right.fill.fore_color.rgb = HEADER_RCB_RED
    nav_right.text_frame.text = "BY:Er. Rk Gupta"
    nav_right.text_frame.paragraphs[0].font.size = HEADER_FONT_SIZE
    nav_right.text_frame.paragraphs[0].font.bold = True
    nav_right.text_frame.paragraphs[0].font.color.rgb = HEADER_YELLOW
    nav_right.text_frame.paragraphs[0].font.name = FONT_NAME
    nav_right.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Determine question box dimensions
    max_width = prs.slide_width.inches * 72 - Inches(1).pt  # Convert slide width to points, leaving margin
    font_size = MAX_FONT_SIZE

    wrapped_question = wrap_text_to_fit(question, font_size=Pt(20), max_width=600)
    question_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(13), Inches(1.5))

    question_frame = question_box.text_frame
    question_frame.clear()
    question_paragraph = question_frame.add_paragraph()
    question_paragraph.text = f"{question_number}. {wrapped_question}"
    question_paragraph.font.bold = True
    question_paragraph.font.color.rgb = QUESTION_TEXT_COLOR
    question_paragraph.font.name = FONT_NAME
    question_paragraph.alignment = PP_ALIGN.LEFT

    # Add year just below the question
    if year:
        year_box = slide.shapes.add_textbox(question_box.left, question_box.top + question_box.height, question_box.width, Inches(0.4))
        year_frame = year_box.text_frame
        year_frame.clear()
        year_paragraph = year_frame.add_paragraph()
        year_paragraph.text = f"{year}"
        year_paragraph.font.size = HEADER_FONT_SIZE
        year_paragraph.font.color.rgb = HEADER_YELLOW
        year_paragraph.font.name = FONT_NAME
        year_paragraph.alignment = PP_ALIGN.RIGHT
        options_top = year_box.top + year_box.height
    else:
        options_top = question_box.top + question_box.height  # Adjust position if year is not provided

    # Shift options just below the year or question
    options_box = slide.shapes.add_textbox(question_box.left, options_top, question_box.width, Inches(3))
    options_frame = options_box.text_frame
    options_frame.clear()
    for option in options:
        option_paragraph = options_frame.add_paragraph()
        option_paragraph.text = option
        option_paragraph.font.color.rgb = OPTION_TEXT_COLOR
        option_paragraph.font.name = FONT_NAME
        option_paragraph.alignment = PP_ALIGN.RIGHT  # Right-align the options

# Test data
new_data = {
    "Question": [
       "उच्च न्यायालयों में तदर्थ न्यायाधीशों की नियुक्ति के संबंध में निम्नलिखित कथनों पर विचार करें:\n a) तदर्थ न्यायाधीशों की नियुक्ति केवल तब की जा सकती है जब न्यायिक रिक्तियां 20% से अधिक हों।\n b) ऐसे न्यायाधीश सभी प्रकार के मामलों की सुनवाई कर सकते हैं, जिसमें संवैधानिक मामले भी शामिल हैं।\n c) सुप्रीम कोर्ट ने हाल ही में उच्च न्यायालयों को आपराधिक मामलों के लिए सेवानिवृत्त न्यायाधीशों की नियुक्ति की अनुमति दी है।\nसही उत्तर चुनें:",
        "हाल ही में चर्चा में रहा सेंट्रल एशियन फ्लाईवे किसके लिए महत्वपूर्ण है?",
        "निम्नलिखित में से कौन-से कोष भारत के जलवायु वित्त (Climate Finance) में योगदान देते हैं?\n a) वैश्विक पर्यावरण सुविधा (Global Environment Facility)\n b) ग्रीन क्लाइमेट फंड (Green Climate Fund)\n c) अनुकूलन कोष (Adaptation Fund)\nनीचे दिए गए कूट का प्रयोग करके सही उत्तर चुनें:",
        "श्रीलंका में अडानी पवन ऊर्जा परियोजना विवादास्पद रही है क्योंकि:\n a) इसका सेंट्रल एशियन फ्लाईवे पर प्रभाव पड़ा है।\n b) परियोजना अनुमोदन प्रक्रिया में कथित भ्रष्टाचार के कारण।",
        "भारत की AI नीति के संबंध में निम्नलिखित कथनों पर विचार करें:\n a) IndiaAI Mission स्टार्टअप और अकादमिक संस्थानों के लिए GPU क्लस्टर को सब्सिडी प्रदान करने का लक्ष्य रखता है।\n b) यह सीधे ChatGPT और DeepSeek जैसे मॉडलों से प्रतिस्पर्धा करता है।\n c) इसका फोकस भारत की संप्रभुता के लिए Foundational Models के प्रशिक्षण पर है।\nसही उत्तर चुनें:",
        "निम्नलिखित में से कौन-सा पिनाका मल्टीपल रॉकेट लॉन्च सिस्टम (MRLS) का सही वर्णन करता है?",
        "ईरान में चाबहार बंदरगाह के विकास का भारत के लिए क्या महत्व है?",
        "क्लिनिकल एस्टैब्लिशमेंट (रजिस्ट्रेशन एंड रेगुलेशन) एक्ट का उद्देश्य क्या है?",
        "हाल ही में अमेरिका द्वारा जारी किया गया नेशनल सिक्योरिटी प्रेसिडेंशियल मेमोरेंडम किसे लक्षित करता है?",
        "हाल ही में पर्यावरण संरक्षण नियमों में संशोधन किया गया था:",
        "राज्य विधेयकों को मंजूरी देने से संबंधित राज्यपाल की शक्तियाँ किस अनुच्छेद से प्राप्त होती हैं?",
        "अमेरिका से अवैध भारतीय प्रवासियों के निर्वासन को भारतीय सरकार ने किस आधार पर उचित ठहराया है?",
        "भारत के सर्वोच्च न्यायालय का हालिया निर्णय तदर्थ (ad-hoc) न्यायाधीशों पर केंद्रित था:",
        "श्रीलंका में अदानी पवन ऊर्जा परियोजना (Adani Wind Power Project) के विरोध का मुख्य कारण क्या है?",
        "‘Flue Gas Desulfurization’ (FGD) तकनीक किसके नियंत्रण में मदद करती है?"
    
    ],
    "Options": [
       ["a और b", "b और c", "केवल c", "a, b और c"],
        ["a) मध्य एशिया और दक्षिण एशिया के बीच व्यापार मार्ग", "b) जलपक्षी प्रजातियों का प्रवास", "c) मध्य एशिया और दक्षिण एशिया के बीच सांस्कृतिक आदान-प्रदान", "d) मानसूनी हवाओं की गति"],
        ["a) केवल a और b", "b) केवल b और c", "c) केवल a और c", "d) a, b और c"],
        ["a) केवल a", "b) केवल b", "c) a और b दोनों", "d) न तो a और न ही b"],
        ["a) केवल a", "b) केवल a और c", "c) केवल b और c", "d) a, b और c"],
        ["a) एक मिसाइल रक्षा प्रणाली", "b) एक लंबी दूरी की तोपखाने रॉकेट प्रणाली", "c) एक नौसैनिक एंटी-शिप मिसाइल", "d) एक एंटी-एयरक्राफ्ट सिस्टम"],
        ["a) यह पाकिस्तान को बाईपास करके अफगानिस्तान के लिए एक वैकल्पिक व्यापार मार्ग प्रदान करता है।", "b) यह ईरानी तेल के आयात से भारत की ऊर्जा सुरक्षा को मजबूत करता है।", "c) यह भारत और ईरान के बीच सैन्य सहयोग को बढ़ाता है।", "d) यह भारत का पहला विदेशी नौसैनिक अड्डा है।"],
        ["a) निजी अस्पतालों और डायग्नोस्टिक लैब को विनियमित करना।", "b) मुफ्त स्वास्थ्य सेवाएं प्रदान करना।", "c) चिकित्सा अनुसंधान गतिविधियों को नियंत्रित करना।", "d) ग्रामीण क्षेत्रों में मेडिकल कॉलेज स्थापित करना।"],
        ["a) चीन की बेल्ट एंड रोड इनिशिएटिव", "b) ईरान के चाबहार बंदरगाह परियोजना", "c) यूक्रेन में रूस की भागीदारी", "d) उत्तर कोरिया का परमाणु कार्यक्रम"],
        ["a) कोयला आधारित विद्युत संयंत्रों के लिए उत्सर्जन मानदंडों को मजबूत करना", "b) सल्फर डाइऑक्साइड मानकों के अनुपालन की समय सीमा बढ़ाना", "c) उद्योगों में सिंगल-यूज़ प्लास्टिक के उपयोग पर प्रतिबंध लगाना", "d) उद्योगों के लिए सख्त जल संरक्षण उपायों को लागू करना"],
        ["a) अनुच्छेद 200", "b) अनुच्छेद 356", "c) अनुच्छेद 239", "d) अनुच्छेद 123"],
        ["a) राजनयिक संबंधों पर विएना सम्मेलन", "b) अमेरिकी आव्रजन और सीमा शुल्क प्रवर्तन (ICE) की मानक संचालन प्रक्रिया", "c) प्रवासन नियंत्रण पर एक द्विपक्षीय संधि", "d) अवैध प्रवासन पर संयुक्त राष्ट्र का प्रस्ताव"],
        ["a) आपराधिक मामलों के लंबित मामलों को कम करने के लिए", "b) न्यायाधीशों की सेवानिवृत्ति आयु बढ़ाने के लिए", "c) सेवानिवृत्त न्यायाधीशों को सभी प्रकार के मामलों की अध्यक्षता करने की अनुमति देने के लिए", "d) भारतीय उच्च न्यायालयों में विदेशी न्यायाधीशों की नियुक्ति के लिए"],
        ["a) परियोजना में चीनी (China) भागीदारी को लेकर चिंताएं", "b) प्रवासी पक्षियों और मछुआरों की आजीविका पर प्रभाव", "c) परमाणु प्रदूषण का खतरा", "d) श्रीलंका की घरेलू राजनीति से इसका संबंध"],
        ["a) ग्रीनहाउस गैस उत्सर्जन", "b) कोयला संयंत्रों से सल्फर डाइऑक्साइड उत्सर्जन", "c) ओजोन परत का क्षरण", "d) उद्योगों से जल प्रदूषण"]
    
    
        ],
   
    "Answers": [
        
    ],
    "Explanation": [
    ]
}

if __name__ == "__main__":
    create_custom_presentation(
    data=new_data,
    title="NEXT LEVEL ACADEMY",
    topic="Daily Current Affairs",
    teacher_name="Prince Sir")