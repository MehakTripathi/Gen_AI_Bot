# FONT_NAME = "Arial"
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Constants for styling

HEADER_TEXT_COLOR = RGBColor(255, 255, 255)  # White
ANSWER_TEXT_COLOR = RGBColor(0, 255, 0)  # Green for answers
FONT_NAME = "Calibri"
HEADER_FONT_SIZE = Pt(16)
QUESTION_FONT_SIZE = Pt(20)
OPTION_FONT_SIZE = Pt(18)
ANSWER_FONT_SIZE = Pt(22)
BACKGROUND_COLOR = RGBColor(0, 0, 0)  # Black
HEADER_RCB_RED = RGBColor(252, 5, 5)
HEADER_YELLOW = RGBColor(255, 255, 0)  # Yellow
HEADER_TEXT_WHITE = RGBColor(255, 255, 255)
QUESTION_TEXT_COLOR = RGBColor(255, 255, 0)  # Yellow
OPTION_TEXT_COLOR = RGBColor(255, 255, 255)  # White
FONT_NAME = "Calibri"
HEADER_FONT_SIZE = Pt(16)
MAX_FONT_SIZE = Pt(28)
MIN_FONT_SIZE = Pt(12)


IMAGE_1_PATH = "logo.jpg"
class PPTHandler:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(14)
        self.prs.slide_height = Inches(7.5)

    @staticmethod
    def _calculate_dynamic_font_size(text, max_font_size=Pt(36), min_font_size=Pt(16)):
        """
        Dynamically calculate font size based on text length.
        
        :param text: Input text
        :param max_font_size: Maximum font size
        :param min_font_size: Minimum font size
        :return: Appropriate font size
        """
        # Basic length-based font size calculation
        length = len(text)
        
        # Adjust font size based on text length
        if length < 50:
            return max_font_size  # Large font for short questions
        elif length < 100:
            return Pt(30)
        elif length < 200:
            return Pt(24)
        elif length < 300:
            return Pt(20)
        else:
            return min_font_size

    @staticmethod
    def wrap_text_to_fit(text, font_size=Pt(36), max_width=700):
        """
        Wrap text to fit within specified width, with dynamic font sizing.
        
        :param text: Input text
        :param font_size: Base font size 
        :param max_width: Maximum width in points
        :return: Wrapped text with adjusted dimensions
        """
        # Dynamically adjust font size based on text length
        adjusted_font_size = PPTHandler._calculate_dynamic_font_size(text)
        
        words = text.split()
        lines = []
        current_line = ""
        approx_char_width = adjusted_font_size.pt * 0.4  # Approximation of character width in points
        
        for word in words:
            if (len(current_line) + len(word)) * approx_char_width > max_width:
                lines.append(current_line.strip())
                current_line = word
            else:
                current_line += " " + word if current_line else word
        
        if current_line:
            lines.append(current_line.strip())

        # Calculate line height with some extra spacing
        line_height = adjusted_font_size.pt * 1.2
        total_height = len(lines) * line_height

        # Return wrapped text and adjust height calculation
        return "\n".join(lines), total_height

    def add_question_slide(prs, question_number, question, options, title, topic, teacher_name):
        """Add a slide for the question with improved sizing and positioning."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Set slide background color
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BACKGROUND_COLOR

        # Add the header with title, logo, and name
        PPTHandler.add_navigation_bar(slide, title, topic, teacher_name)

        # Determine question box dimensions dynamically
        font_size = Pt(36)
        wrapped_question, text_height = PPTHandler.wrap_text_to_fit(question, font_size=font_size, max_width=700)
        
        # Adjust box dimensions based on text height
        box_width = Inches(6.3)  # Default width for the box
        lines = wrapped_question.count('\n') + 1 
        
        # Calculate dynamic box height with padding
        box_height = Inches((text_height / 72) + 0.2)  # Convert points to inches, with padding

        # Create question box
        question_box = slide.shapes.add_textbox(Inches(7.5), Inches(1), box_width, box_height)

        # Add question number box
        question_number_box = slide.shapes.add_textbox(Inches(6.3), Inches(1.3), Inches(0.7), Inches(0.9))
        question_number_frame = question_number_box.text_frame
        question_number_frame.clear()
        question_number_paragraph = question_number_frame.add_paragraph()
        question_number_paragraph.text = str(question_number) + "."  # Add the question number
        question_number_box.fill.solid()
        question_number_box.fill.fore_color.rgb = RGBColor(255, 0, 0)
        question_number_paragraph.font.bold = True
        question_number_paragraph.font.color.rgb = RGBColor(255, 255, 255)
        question_number_paragraph.font.name = FONT_NAME
        question_number_paragraph.font.size = Pt(30)
        question_number_paragraph.alignment = PP_ALIGN.RIGHT

        # Populate the question text box
        question_frame = question_box.text_frame
        question_frame.clear()
        question_frame.word_wrap = True
        question_frame.margin_top = Pt(0)
        question_frame.margin_bottom = Pt(0)
        question_frame.margin_left = Pt(0)
        question_frame.margin_right = Pt(0)
        question_paragraph = question_frame.add_paragraph()
        question_paragraph.text = wrapped_question
        question_paragraph.font.bold = True
        question_paragraph.font.color.rgb = QUESTION_TEXT_COLOR
        question_paragraph.font.name = FONT_NAME
        
        # Dynamically adjust font size
        dynamic_font_size = PPTHandler._calculate_dynamic_font_size(question)
        question_paragraph.font.size = dynamic_font_size
        question_paragraph.alignment = PP_ALIGN.LEFT

    
        total_options_height = 0
        wrapped_options = []
        
        # Process each option for wrapping and height calculation
        for option in options:
            wrapped_option, option_height = PPTHandler.wrap_text_to_fit(
                option, 
                font_size=OPTION_FONT_SIZE,
                max_width=600  # Slightly smaller than box width to ensure padding
            )
            wrapped_options.append(wrapped_option)
            total_options_height += option_height + Pt(10).pt  # Add padding between options
        
        # Calculate optimal position for options box
        # Start options 1 inch below question box, or at slide middle, whichever is higher
        min_options_top = question_box.top + question_box.height + Inches(1)
        slide_middle = (prs.slide_height - Inches(1)) / 2  # Leave space for header
        options_top = max(min_options_top, slide_middle - Inches(total_options_height/72/2))
        
        # Create options box with calculated height
        options_box = slide.shapes.add_textbox(
            question_box.left, 
            options_top, 
            question_box.width, 
            Inches(total_options_height/72 + 0.5)  # Convert points to inches, add padding
        )
        
        options_frame = options_box.text_frame
        options_frame.clear()
        options_frame.word_wrap = True
        
        # Add each wrapped option
        for wrapped_option in wrapped_options:
            option_paragraph = options_frame.add_paragraph()
            option_paragraph.text = wrapped_option
            option_paragraph.font.size = OPTION_FONT_SIZE
            option_paragraph.line_spacing = 1.2  # Proportional line spacing
            option_paragraph.space_after = Pt(10)  # Space between options
            option_paragraph.font.color.rgb = OPTION_TEXT_COLOR
            option_paragraph.font.name = FONT_NAME
            option_paragraph.alignment = PP_ALIGN.LEFT

    def create_custom_presentation(self, data, title, topic, teacher_name):
        """Create a custom PowerPoint presentation based on the provided data."""
        print("Generating New PPT...")

        # Create a new PowerPoint presentation
        prs = Presentation()
        prs.slide_width = Inches(14)
        prs.slide_height = Inches(7.5)

        question_count = len(data["Question"])
        year_count = len(data.get("Year", []))
        options_count = len(data.get("Options", []))

        if question_count == 0:
            print("No data to generate PPT.")
            return None  # Avoid generating an empty PPT

        for i in range(question_count):
            question = data["Question"][i]
            year = data["Year"][i] if i < year_count else ""
            options = data["Options"][i] if i < options_count else []

            # Debugging prints to check data
            print(f"Processing Question {i+1}: {question}")
            print(f"Options: {options}")

            # Add Question Slide
            PPTHandler.add_question_slide(prs, i + 1, question, options, title, topic, teacher_name)

        # Ensure unique filename every time
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"ppt_{timestamp}.pptx"

        prs.save(filename)
        print(f"Presentation saved as {filename}")

        return filename  # Return the filename to be sent in Telegram bot


    def add_navigation_bar(slide,title,topic,teacher_name):
        """Add the navigation bar with title, logo, and name."""
        # Title Section
        nav_left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(0), Inches(8), Inches(0.5))
        nav_left.fill.solid()
        nav_left.fill.fore_color.rgb = HEADER_RCB_RED
        nav_left.text_frame.text = title
        nav_left.text_frame.paragraphs[0].font.size = HEADER_FONT_SIZE
        nav_left.text_frame.paragraphs[0].font.bold = True
        nav_left.text_frame.paragraphs[0].font.color.rgb = HEADER_TEXT_WHITE
        nav_left.text_frame.paragraphs[0].font.name = FONT_NAME
        nav_left.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Add logo
        slide.shapes.add_picture(IMAGE_1_PATH, Inches(0), Inches(0), Inches(1), Inches(0.5))

        # Middle Section for Topic
        nav_middle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(0), Inches(6), Inches(0.5))
        nav_middle.fill.solid()
        nav_middle.fill.fore_color.rgb = HEADER_YELLOW
        nav_middle.text_frame.clear()  # Clear any default text
        nav_middle.text_frame.text = topic  # Add your topic here
        nav_middle.text_frame.paragraphs[0].font.size = HEADER_FONT_SIZE
        nav_middle.text_frame.paragraphs[0].font.bold = True
        nav_middle.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        nav_middle.text_frame.paragraphs[0].font.name = FONT_NAME
        nav_middle.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Right Section for Creator's Name
        nav_right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10), Inches(0), Inches(4), Inches(0.5))
        nav_right.fill.solid()
        nav_right.fill.fore_color.rgb = HEADER_RCB_RED
        nav_right.text_frame.text = f"BY: {teacher_name}"
        nav_right.text_frame.paragraphs[0].font.size = HEADER_FONT_SIZE
        nav_right.text_frame.paragraphs[0].font.bold = True
        nav_right.text_frame.paragraphs[0].font.color.rgb = HEADER_TEXT_WHITE
        nav_right.text_frame.paragraphs[0].font.name = FONT_NAME
        nav_right.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE